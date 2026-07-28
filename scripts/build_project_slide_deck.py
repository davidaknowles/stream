#!/usr/bin/env python
"""Build a slide deck summarizing STREAM model ideas and results."""

from __future__ import annotations

import argparse
import tempfile
from pathlib import Path

import pandas as pd
from PIL import Image
from plotnine import (
    aes,
    coord_cartesian,
    element_blank,
    element_text,
    facet_wrap,
    geom_col,
    geom_label,
    geom_text,
    geom_tile,
    ggplot,
    labs,
    position_dodge,
    scale_fill_brewer,
    scale_fill_gradient,
    scale_fill_manual,
    theme,
    theme_bw,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]

BG = RGBColor(248, 249, 250)
INK = RGBColor(32, 37, 42)
MUTED = RGBColor(92, 103, 115)
BLUE = RGBColor(49, 101, 163)
TEAL = RGBColor(43, 145, 134)
GREEN = RGBColor(83, 154, 93)
GOLD = RGBColor(214, 158, 46)
RED = RGBColor(190, 75, 73)
PURPLE = RGBColor(112, 93, 180)
LIGHT_BLUE = RGBColor(225, 235, 247)
LIGHT_GREEN = RGBColor(224, 241, 229)
LIGHT_GOLD = RGBColor(249, 238, 207)
LIGHT_RED = RGBColor(249, 224, 220)
WHITE = RGBColor(255, 255, 255)

SLIDE_W = 13.333
SLIDE_H = 7.5


def inches(value: float):
    return Inches(value)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
) -> None:
    box = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = inches(0.04)
    frame.margin_right = inches(0.04)
    frame.margin_top = inches(0.02)
    frame.margin_bottom = inches(0.02)
    frame.vertical_anchor = valign
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.55, 0.26, 12.2, 0.48, size=25, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.57, 0.79, 11.7, 0.28, size=11, color=MUTED)


def set_background(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = RGBColor(210, 216, 222),
    size: int = 13,
    bold: bool = False,
    color: RGBColor = INK,
    radius: bool = True,
):
    shape_type = 5 if radius else 1
    shape = slide.shapes.add_shape(shape_type, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.1)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = inches(0.08)
    frame.margin_right = inches(0.08)
    frame.margin_top = inches(0.06)
    frame.margin_bottom = inches(0.06)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = MUTED, width: float = 1.5) -> None:
    line = slide.shapes.add_connector(1, inches(x1), inches(y1), inches(x2), inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)


def add_metric(slide, label: str, value: str, x: float, y: float, w: float, h: float, fill: RGBColor) -> None:
    add_box(slide, "", x, y, w, h, fill=fill, line=RGBColor(210, 216, 222), radius=False)
    add_text(slide, value, x + 0.06, y + 0.14, w - 0.12, 0.36, size=23, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.08, y + 0.58, w - 0.16, 0.4, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    draw_w = iw * scale
    draw_h = ih * scale
    slide.shapes.add_picture(str(path), inches(x + (w - draw_w) / 2), inches(y + (h - draw_h) / 2), inches(draw_w), inches(draw_h))


def add_table(slide, frame: pd.DataFrame, x: float, y: float, w: float, h: float, *, font_size: int = 10) -> None:
    rows = len(frame) + 1
    cols = len(frame.columns)
    table = slide.shapes.add_table(rows, cols, inches(x), inches(y), inches(w), inches(h)).table
    for col_idx, column in enumerate(frame.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
    for row_idx, row in enumerate(frame.itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else RGBColor(241, 244, 247)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.name = "Aptos"
            run.font.size = Pt(font_size)
            run.font.color.rgb = INK


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    return slide


def collect_mouse_summary() -> pd.DataFrame:
    rows = []
    label_map = {
        "standard_cfm": "CFM",
        "film": "FiLM",
        "cross_attention": "Cross-attn",
        "standard_cfm_uce": "CFM + UCE",
        "film_uce": "FiLM + UCE",
        "cross_attention_uce": "Cross-attn + UCE",
    }
    for panel in (5000, 10000):
        for path in sorted((ROOT / f"outputs/stream_hvg{panel}").glob("eval_metrics_*.csv")):
            model = path.stem.removeprefix("eval_metrics_")
            frame = pd.read_csv(path)
            frame["panel"] = panel
            frame["model"] = model
            frame["model_label"] = label_map[model]
            frame["state"] = "UCE state" if model.endswith("_uce") else "Expression state"
            frame["family"] = model.removesuffix("_uce")
            rows.append(frame)
    combined = pd.concat(rows, ignore_index=True)
    return (
        combined.groupby(["panel", "model", "model_label", "state", "family", "eval_gene_set", "n_eval_genes"], as_index=False)
        .agg(mean_loss=("loss", "mean"), mean_mae=("velocity_mae", "mean"))
    )


def collect_original_legacy() -> pd.DataFrame:
    rows = []
    label_map = {"standard_cfm": "CFM", "film": "FiLM", "cross_attention": "Cross-attn"}
    for path in sorted((ROOT / "outputs/stream").glob("eval_metrics_*.csv")):
        model = path.stem.removeprefix("eval_metrics_")
        frame = pd.read_csv(path)
        rows.append(
            {
                "family": model,
                "model_label": label_map[model],
                "trained_panel": "1,984",
                "mean_loss": frame["loss"].mean(),
                "mean_mae": frame["velocity_mae"].mean(),
            }
        )
    return pd.DataFrame(rows)


def collect_legacy_panel_summary(mouse_summary: pd.DataFrame) -> pd.DataFrame:
    label_map = {"standard_cfm": "CFM", "film": "FiLM", "cross_attention": "Cross-attn"}
    rows = [collect_original_legacy()]
    for panel in (5000, 10000):
        subset = mouse_summary[(mouse_summary["eval_gene_set"] == "legacy_1984") & (mouse_summary["state"] == "Expression state")].copy()
        subset = subset[subset["panel"] == panel]
        rows.append(
            pd.DataFrame(
                {
                    "family": subset["family"],
                    "model_label": subset["family"].map(label_map),
                    "trained_panel": f"{panel:,}",
                    "mean_loss": subset["mean_loss"],
                    "mean_mae": subset["mean_mae"],
                }
            )
        )
    uce = mouse_summary[
        (mouse_summary["panel"] == 10000)
        & (mouse_summary["eval_gene_set"] == "legacy_1984")
        & (mouse_summary["state"] == "UCE state")
    ].copy()
    rows.append(
        pd.DataFrame(
            {
                "family": uce["family"],
                "model_label": uce["family"].map(label_map),
                "trained_panel": "10,000 + UCE",
                "mean_loss": uce["mean_loss"],
                "mean_mae": uce["mean_mae"],
            }
        )
    )
    out = pd.concat(rows, ignore_index=True)
    order = {"1,984": 0, "5,000": 1, "10,000": 2, "10,000 + UCE": 3}
    out["panel_order"] = out["trained_panel"].map(order)
    return out.sort_values(["model_label", "panel_order"])


def collect_zebrafish_summary() -> pd.DataFrame:
    path = ROOT / "figures/zebrafish_transfer_displacement_mae.csv"
    if not path.exists():
        raise FileNotFoundError(path)
    frame = pd.read_csv(path)
    frame = frame[frame["panel"] == 10000].copy()
    frame = frame[frame["training"] != "no_model"].copy()
    frame["training_label"] = frame["training"].map(
        {
            "zero_shot": "Mouse zero-shot",
            "fine_tuned": "Mouse -> zebrafish",
            "zebrafish_only": "Zebrafish only",
        }
    )
    frame["time_label"] = frame["time_scale"].map({"relative": "Relative time", "days": "Physical days"})
    return frame


def dataset_summary() -> pd.DataFrame:
    mouse = pd.read_csv(ROOT / "outputs/jax_adata_eda/adata_summary.csv")
    z_manifest = pd.read_csv(ROOT / "outputs/zscape_uce_full_controls/embeddings/manifest.csv")
    return pd.DataFrame(
        [
            {
                "Species": "Mouse",
                "Cells": f"{mouse['n_obs'].sum() / 1e6:.2f}M",
                "Timepoints": "43",
                "Held out": "E9.5, E10.5",
                "Genes": "5k / 10k",
                "CRE tokens": "318k",
            },
            {
                "Species": "Zebrafish",
                "Cells": f"{z_manifest['cells'].sum() / 1e6:.2f}M",
                "Timepoints": "18",
                "Held out": "36, 72 hpf",
                "Genes": "10k",
                "CRE tokens": "318k",
            },
        ]
    )


def save_mouse_full_plot(summary: pd.DataFrame, path: Path) -> None:
    frame = summary[summary["eval_gene_set"] == "full"].copy()
    frame["panel_label"] = frame["panel"].map({5000: "5k genes", 10000: "10k genes"})
    family_order = ["standard_cfm", "film", "cross_attention"]
    frame["family_label"] = pd.Categorical(
        frame["family"].map({"standard_cfm": "CFM", "film": "FiLM", "cross_attention": "Cross-attn"}),
        categories=["CFM", "FiLM", "Cross-attn"],
        ordered=True,
    )
    plot = (
        ggplot(frame, aes("family_label", "mean_loss", fill="state"))
        + geom_col(position=position_dodge(width=0.72), width=0.65)
        + facet_wrap("~panel_label", scales="free_y")
        + scale_fill_manual(values={"Expression state": "#5B8FD1", "UCE state": "#57A773"})
        + labs(x="", y="Held-out MSE", fill="")
        + theme_bw(base_size=12)
        + theme(
            figure_size=(7.8, 3.9),
            subplots_adjust={"wspace": 0.18},
            panel_grid_minor=element_blank(),
            axis_text_x=element_text(rotation=18, ha="right"),
            legend_position="top",
        )
    )
    plot.save(path, width=7.8, height=3.9, dpi=180, verbose=False)


def save_legacy_heatmap(frame: pd.DataFrame, path: Path) -> None:
    heat = frame.copy()
    heat["loss_label"] = heat["mean_loss"].map(lambda x: f"{x:.1f}")
    heat["trained_panel"] = pd.Categorical(
        heat["trained_panel"],
        categories=["1,984", "5,000", "10,000", "10,000 + UCE"],
        ordered=True,
    )
    heat["model_label"] = pd.Categorical(heat["model_label"], categories=["CFM", "FiLM", "Cross-attn"], ordered=True)
    plot = (
        ggplot(heat, aes("trained_panel", "model_label", fill="mean_loss"))
        + geom_tile(color="white", size=1.3)
        + geom_text(aes(label="loss_label"), size=11, color="#1f252a")
        + scale_fill_gradient(low="#DDF0DF", high="#F3B3A8")
        + labs(x="Modeled gene panel", y="", fill="MSE")
        + theme_bw(base_size=12)
        + theme(
            figure_size=(7.8, 3.25),
            panel_grid_major=element_blank(),
            panel_grid_minor=element_blank(),
            axis_text_x=element_text(rotation=0),
            legend_position="right",
        )
    )
    plot.save(path, width=7.8, height=3.25, dpi=180, verbose=False)


def save_zebrafish_plot(summary: pd.DataFrame, path: Path) -> None:
    frame = summary.copy()
    frame = frame[frame["eval_gene_set"] == "full"].copy()
    frame["training_label"] = pd.Categorical(
        frame["training_label"],
        categories=["Mouse zero-shot", "Mouse -> zebrafish", "Zebrafish only"],
        ordered=True,
    )
    plot = (
        ggplot(frame, aes("training_label", "displacement_mae", fill="training_label"))
        + geom_col(width=0.72)
        + facet_wrap("~time_label")
        + scale_fill_manual(values=["#D89B3D", "#57A773", "#5B8FD1"])
        + coord_cartesian(ylim=(0, 0.32))
        + labs(x="", y="Displacement MAE", fill="")
        + theme_bw(base_size=12)
        + theme(
            figure_size=(7.8, 3.75),
            axis_text_x=element_text(rotation=20, ha="right"),
            legend_position="none",
            panel_grid_minor=element_blank(),
        )
    )
    plot.save(path, width=7.8, height=3.75, dpi=180, verbose=False)


def create_deck(output: Path) -> None:
    mouse_summary = collect_mouse_summary()
    legacy_summary = collect_legacy_panel_summary(mouse_summary)
    zfish_summary = collect_zebrafish_summary()
    data_summary = dataset_summary()

    prs = Presentation()
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)

    with tempfile.TemporaryDirectory() as tmp_raw:
        tmp = Path(tmp_raw)
        mouse_full_plot = tmp / "mouse_full_panel.png"
        legacy_heatmap = tmp / "legacy_heatmap.png"
        zfish_plot = tmp / "zfish_transfer.png"
        save_mouse_full_plot(mouse_summary, mouse_full_plot)
        save_legacy_heatmap(legacy_summary, legacy_heatmap)
        save_zebrafish_plot(zfish_summary, zfish_plot)

        # 1
        slide = blank_slide(prs)
        add_text(slide, "STREAM learns developmental vector fields from cell state and regulatory sequence", 0.55, 0.35, 9.0, 1.1, size=28, bold=True)
        add_text(slide, "Multi-timepoint CFM + CRE sequence conditioning + UCE cell-state embeddings", 0.58, 1.38, 9.0, 0.35, size=15, color=MUTED)
        add_box(slide, "Time-resolved\nsingle-cell atlas", 0.78, 2.45, 2.25, 1.05, fill=LIGHT_BLUE, bold=True, line=BLUE)
        add_box(slide, "Regulatory DNA\ncCRE tokens", 3.55, 2.45, 2.25, 1.05, fill=LIGHT_GOLD, bold=True, line=GOLD)
        add_box(slide, "UCE or expression\ncell state", 6.32, 2.45, 2.25, 1.05, fill=LIGHT_GREEN, bold=True, line=GREEN)
        add_box(slide, "Gene velocity\nfield", 9.1, 2.45, 2.25, 1.05, fill=WHITE, bold=True, line=PURPLE)
        for x1, x2 in [(3.03, 3.55), (5.8, 6.32), (8.57, 9.1)]:
            add_line(slide, x1, 2.98, x2, 2.98, color=MUTED, width=2)
        add_metric(slide, "mouse atlas cells", "11.44M", 0.9, 5.1, 2.2, 1.05, LIGHT_BLUE)
        add_metric(slide, "zebrafish control cells", "1.23M", 3.45, 5.1, 2.2, 1.05, LIGHT_BLUE)
        add_metric(slide, "modeled genes", "10k", 6.0, 5.1, 2.2, 1.05, LIGHT_GOLD)
        add_metric(slide, "best mouse model", "UCE + cross-attn", 8.55, 5.1, 3.0, 1.05, LIGHT_GREEN)

        # 2
        slide = blank_slide(prs)
        add_title(slide, "Development is observed as snapshots, but the target is a continuous vector field")
        add_image_fit(slide, ROOT / "figures/full_umap_by_major_trajectory.png", 0.55, 1.08, 5.9, 5.9)
        add_table(slide, data_summary, 6.7, 1.25, 5.95, 1.65, font_size=9)
        add_box(slide, "Held-out timepoints test temporal interpolation around missing stages", 7.0, 3.45, 5.05, 0.72, fill=LIGHT_RED, line=RED, bold=True)
        add_box(slide, "The same output space is used for CFM and STREAM, so sequence conditioning is the isolated change", 7.0, 4.45, 5.05, 0.82, fill=WHITE, line=RGBColor(210, 216, 222), size=12)
        add_box(slide, "Cross-species tests swap regulatory sequence and UCE gene tokens while retaining expression-valued targets", 7.0, 5.55, 5.05, 0.95, fill=WHITE, line=RGBColor(210, 216, 222), size=12)

        # 3
        slide = blank_slide(prs)
        add_title(slide, "Multi-timepoint CFM trains on adjacent intervals and scores intervals touching held-out stages")
        y = 2.35
        xs = [0.9, 2.05, 3.2, 4.35, 5.5, 6.65, 7.8, 8.95, 10.1, 11.25]
        labels = ["E8.5", "E8.75", "E9.0", "E9.25", "E9.5", "E9.75", "E10.0", "E10.25", "E10.5", "E10.75"]
        held = {"E9.5", "E10.5"}
        for x, lab in zip(xs, labels, strict=True):
            fill = LIGHT_RED if lab in held else WHITE
            line = RED if lab in held else BLUE
            add_box(slide, lab, x, y, 0.82, 0.48, fill=fill, line=line, size=10, bold=lab in held, radius=False)
        for i in range(len(xs) - 1):
            color = RED if labels[i] in held or labels[i + 1] in held else GREEN
            add_line(slide, xs[i] + 0.82, y + 0.24, xs[i + 1], y + 0.24, color=color, width=2.2)
        add_text(slide, "green = trained adjacent intervals     red = held-out-touching evaluation intervals", 2.3, 3.05, 8.6, 0.25, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        add_box(slide, "1. sample cells at t_k", 0.9, 4.25, 2.2, 0.7, fill=LIGHT_BLUE, line=BLUE, bold=True)
        add_box(slide, "2. minibatch OT pairs cells", 3.45, 4.25, 2.3, 0.7, fill=LIGHT_GOLD, line=GOLD, bold=True)
        add_box(slide, "3. interpolate x_t at random tau", 6.1, 4.25, 2.45, 0.7, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "4. regress velocity (x_1 - x_0) / dt", 8.9, 4.25, 3.0, 0.7, fill=WHITE, line=PURPLE, bold=True)
        for x1, x2 in [(3.1, 3.45), (5.75, 6.1), (8.55, 8.9)]:
            add_line(slide, x1, 4.6, x2, 4.6, color=MUTED, width=2)
        add_text(slide, "One vector field is fit across all training intervals; no ODE rollout is needed during fitting.", 1.05, 5.75, 11.0, 0.45, size=16, bold=True, align=PP_ALIGN.CENTER)

        # 4
        slide = blank_slide(prs)
        add_title(slide, "STREAM predicts each gene velocity from global cell state and local regulatory tokens")
        add_box(slide, "cell state\nx_t or u_t", 0.85, 2.9, 1.8, 0.95, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "cCRE + promoter\ntokens for gene g", 0.85, 1.35, 1.8, 0.95, fill=LIGHT_GOLD, line=GOLD, bold=True)
        add_box(slide, "AlphaGenome\nsequence embeddings", 3.05, 1.35, 2.05, 0.95, fill=WHITE, line=GOLD, bold=True)
        add_box(slide, "shared CRE\ntransformer layers", 5.55, 1.35, 2.15, 2.5, fill=LIGHT_BLUE, line=BLUE, bold=True)
        add_box(slide, "FiLM or\ncross-attention", 8.25, 2.9, 2.1, 0.95, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "promoter token\nreadout", 8.25, 1.35, 2.1, 0.95, fill=WHITE, line=PURPLE, bold=True)
        add_box(slide, "velocity\nf_g(s_t, Z_g)", 10.9, 1.9, 1.55, 1.25, fill=WHITE, line=RED, bold=True)
        for x1, y1, x2, y2 in [
            (2.65, 1.82, 3.05, 1.82),
            (5.1, 1.82, 5.55, 1.82),
            (7.7, 1.82, 8.25, 1.82),
            (10.35, 1.82, 10.9, 2.35),
            (2.65, 3.38, 8.25, 3.38),
            (10.35, 3.38, 10.9, 2.75),
        ]:
            add_line(slide, x1, y1, x2, y2, color=MUTED, width=2)
        add_text(slide, "Gene specificity comes from linked sequence tokens; weights are shared across genes.", 1.4, 5.25, 10.6, 0.4, size=16, bold=True, align=PP_ALIGN.CENTER)

        # 5
        slide = blank_slide(prs)
        add_title(slide, "UCE changes the state representation; CFM targets remain expression-valued")
        add_box(slide, "expression endpoints\nx_k, x_k+1", 0.85, 1.35, 2.0, 0.9, fill=LIGHT_BLUE, line=BLUE, bold=True)
        add_box(slide, "OT pairs and\nlinear paths", 3.2, 1.35, 2.0, 0.9, fill=LIGHT_GOLD, line=GOLD, bold=True)
        add_box(slide, "target expression\nvelocity", 5.55, 1.35, 2.0, 0.9, fill=WHITE, line=RED, bold=True)
        add_box(slide, "frozen UCE\nencoder", 3.2, 3.35, 2.0, 0.9, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "interpolated\nu_t in R^1280", 5.55, 3.35, 2.0, 0.9, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "STREAM / CFM\nvector field", 7.9, 2.35, 2.0, 0.9, fill=WHITE, line=PURPLE, bold=True)
        add_box(slide, "predicted expression\nvelocity", 10.25, 2.35, 2.0, 0.9, fill=WHITE, line=RED, bold=True)
        for x1, y1, x2, y2 in [(2.85, 1.8, 3.2, 1.8), (5.2, 1.8, 5.55, 1.8), (2.85, 1.95, 3.2, 3.8), (5.2, 3.8, 5.55, 3.8), (7.55, 3.8, 7.9, 2.8), (9.9, 2.8, 10.25, 2.8)]:
            add_line(slide, x1, y1, x2, y2, color=MUTED, width=2)
        add_table(
            slide,
            pd.DataFrame(
                [
                    ["Expression state", "x_t in R^G", "G-gene velocity"],
                    ["UCE state", "u_t in R^1280", "G-gene velocity"],
                ],
                columns=["Run", "Input to f", "Output / loss"],
            ),
            1.5,
            5.25,
            10.4,
            1.05,
            font_size=11,
        )

        # 6
        slide = blank_slide(prs)
        add_title(slide, "Mouse: 10k genes lowers full-panel loss; cross-attention is best by MSE")
        add_image_fit(slide, mouse_full_plot, 0.75, 1.15, 8.4, 4.55)
        full10 = mouse_summary[(mouse_summary["panel"] == 10000) & (mouse_summary["eval_gene_set"] == "full")]
        best = full10.sort_values("mean_loss").iloc[0]
        best_mae = full10.sort_values("mean_mae").iloc[0]
        add_metric(slide, "best 10k MSE", f"{best['mean_loss']:.2f}", 9.45, 1.45, 2.4, 1.0, LIGHT_GREEN)
        add_text(slide, str(best["model_label"]), 9.55, 2.42, 2.2, 0.28, size=13, bold=True, align=PP_ALIGN.CENTER)
        add_metric(slide, "best 10k MAE", f"{best_mae['mean_mae']:.3f}", 9.45, 3.05, 2.4, 1.0, LIGHT_BLUE)
        add_text(slide, str(best_mae["model_label"]), 9.55, 4.02, 2.2, 0.28, size=13, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, "Full-panel losses compare each model on its own selected genes; the legacy-panel slide controls for output genes.", 0.9, 6.25, 11.2, 0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)

        # 7
        slide = blank_slide(prs)
        add_title(slide, "Fair legacy-panel scoring shows the 10k cross-attention gain is real")
        add_image_fit(slide, legacy_heatmap, 0.85, 1.2, 8.4, 3.8)
        add_box(slide, "5k models do not improve the 1,984-gene metric", 9.65, 1.5, 2.55, 0.85, fill=LIGHT_RED, line=RED, bold=True)
        add_box(slide, "10k cross-attention improves loss from 84.8 to 80.5", 9.65, 2.75, 2.55, 0.95, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "10k UCE cross-attention improves further to 78.7", 9.65, 4.2, 2.55, 0.95, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_text(slide, "Numbers are held-out MSE on the same original gene panel; lower is better.", 1.0, 6.0, 11.0, 0.4, size=13, color=MUTED, align=PP_ALIGN.CENTER)

        # 8
        slide = blank_slide(prs)
        add_title(slide, "The best mouse model is 10k UCE cross-attention")
        table = mouse_summary[
            (mouse_summary["panel"] == 10000)
            & (mouse_summary["eval_gene_set"] == "full")
            & (mouse_summary["model"].isin(["standard_cfm", "film", "cross_attention", "cross_attention_uce"]))
        ].copy()
        table = table.assign(
            Model=table["model_label"],
            MSE=table["mean_loss"].map(lambda x: f"{x:.2f}"),
            MAE=table["mean_mae"].map(lambda x: f"{x:.3f}"),
        )[["Model", "MSE", "MAE"]].sort_values("MSE")
        add_table(slide, table, 0.9, 1.35, 5.4, 2.2, font_size=13)
        add_box(slide, "UCE helps when the model can use sequence-specific cell context", 7.0, 1.35, 4.7, 0.9, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "Cross-attention can query UCE context differently for each gene's CRE tokens", 7.0, 2.65, 4.7, 1.0, fill=WHITE, line=PURPLE, bold=True)
        add_box(slide, "Expression-state CFM remains a strong baseline, so sequence conditioning must beat a high bar", 7.0, 4.05, 4.7, 1.0, fill=WHITE, line=BLUE, bold=True)
        add_metric(slide, "best full-panel MSE", "20.28", 1.2, 4.65, 2.6, 1.1, LIGHT_GREEN)
        add_metric(slide, "best full-panel MAE", "1.594", 4.05, 4.65, 2.6, 1.1, LIGHT_GREEN)

        # 9
        slide = blank_slide(prs)
        add_title(slide, "Cross-species transfer swaps the genome and cell-state vocabulary, not the training objective")
        add_box(slide, "mouse 10k\nUCE cross-attn", 0.85, 1.45, 2.2, 0.95, fill=LIGHT_BLUE, line=BLUE, bold=True)
        add_box(slide, "source checkpoint", 3.55, 1.45, 1.8, 0.95, fill=WHITE, line=BLUE, bold=True)
        add_box(slide, "zebrafish UCE\ncell states", 0.85, 3.35, 2.2, 0.95, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "zebrafish cCRE\nsequence tokens", 3.55, 3.35, 1.8, 0.95, fill=LIGHT_GOLD, line=GOLD, bold=True)
        add_box(slide, "zero-shot", 6.0, 1.05, 1.7, 0.7, fill=LIGHT_GOLD, line=GOLD, bold=True)
        add_box(slide, "fine-tuned", 6.0, 2.15, 1.7, 0.7, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "zebrafish only", 6.0, 3.25, 1.7, 0.7, fill=LIGHT_BLUE, line=BLUE, bold=True)
        add_box(slide, "held-out 36 and 72 hpf\nexpression displacement", 8.35, 2.45, 3.3, 1.05, fill=WHITE, line=RED, bold=True)
        for x1, y1, x2, y2 in [(3.05, 1.92, 3.55, 1.92), (5.35, 1.92, 6.0, 1.4), (5.35, 3.82, 6.0, 2.5), (5.35, 3.82, 6.0, 3.6), (7.7, 2.55, 8.35, 2.95)]:
            add_line(slide, x1, y1, x2, y2, color=MUTED, width=2)
        add_text(slide, "Two time conventions are evaluated: physical days and organism-relative [0,1] time.", 1.1, 5.7, 10.9, 0.45, size=15, bold=True, align=PP_ALIGN.CENTER)

        # 10
        slide = blank_slide(prs)
        add_title(slide, "Zebrafish transfer is challenging: fine-tuning is best learned model in relative time")
        add_image_fit(slide, zfish_plot, 0.85, 1.15, 8.6, 4.45)
        rel = zfish_summary[(zfish_summary["time_scale"] == "relative") & (zfish_summary["eval_gene_set"] == "full")]
        days = zfish_summary[(zfish_summary["time_scale"] == "days") & (zfish_summary["eval_gene_set"] == "full")]
        add_metric(slide, "relative fine-tuned", f"{rel[rel['training'] == 'fine_tuned']['displacement_mae'].iloc[0]:.3f}", 9.8, 1.35, 2.3, 0.95, LIGHT_GREEN)
        add_metric(slide, "relative zebrafish only", f"{rel[rel['training'] == 'zebrafish_only']['displacement_mae'].iloc[0]:.3f}", 9.8, 2.7, 2.3, 0.95, LIGHT_BLUE)
        add_metric(slide, "physical zero-shot", f"{days[days['training'] == 'zero_shot']['displacement_mae'].iloc[0]:.3f}", 9.8, 4.05, 2.3, 0.95, LIGHT_GOLD)
        add_text(slide, "Lower displacement MAE is better. Bars compare learned transfer and target-only regimes.", 1.0, 6.2, 11.2, 0.35, size=12, color=MUTED, align=PP_ALIGN.CENTER)

        # 11
        slide = blank_slide(prs)
        add_title(slide, "Cross-species generalization depends on the time coordinate and training regime")
        ztab = zfish_summary[zfish_summary["eval_gene_set"] == "full"].copy()
        ztab = ztab.assign(
            Time=ztab["time_label"],
            Regime=ztab["training_label"],
            MAE=ztab["displacement_mae"].map(lambda x: f"{x:.3f}"),
            MSE=ztab["displacement_mse"].map(lambda x: f"{x:.3f}"),
        )[["Time", "Regime", "MAE", "MSE"]]
        ztab["order"] = zfish_summary["time_scale"].map({"relative": 0, "days": 1}).to_numpy()
        ztab["regime_order"] = zfish_summary["training"].map({"zero_shot": 0, "fine_tuned": 1, "zebrafish_only": 2}).to_numpy()
        ztab = ztab.sort_values(["order", "regime_order"]).drop(columns=["order", "regime_order"])
        add_table(slide, ztab, 0.75, 1.15, 7.5, 4.15, font_size=10)
        add_box(slide, "Relative time: fine-tuning beats zero-shot and zebrafish-only among learned models", 8.75, 1.35, 3.4, 1.0, fill=LIGHT_GREEN, line=GREEN, bold=True)
        add_box(slide, "Physical days: zero-shot is the strongest learned model, suggesting transferred timing is useful", 8.75, 2.75, 3.4, 1.05, fill=LIGHT_GOLD, line=GOLD, bold=True)
        add_box(slide, "Zebrafish-only training is weaker than transfer-started models in both time coordinates", 8.75, 4.25, 3.4, 1.05, fill=LIGHT_BLUE, line=BLUE, bold=True)

        # 12
        slide = blank_slide(prs)
        add_title(slide, "Current picture: 10k + UCE + cross-attention works in mouse; transfer needs stronger alignment")
        takeaways = pd.DataFrame(
            [
                ["Mouse panel size", "10k > 5k", "Cross-attn improves legacy-panel loss at 10k"],
                ["State representation", "UCE helps STREAM", "Best mouse model is 10k UCE cross-attn"],
                ["Sequence conditioning", "Useful but not automatic", "CFM remains a strong baseline"],
                ["Cross-species", "Transfer signal appears", "Fine-tuning helps in relative time"],
            ],
            columns=["Question", "Result", "Evidence"],
        )
        add_table(slide, takeaways, 0.75, 1.2, 11.85, 2.75, font_size=12)
        add_box(slide, "Next experiments", 0.95, 4.6, 2.1, 0.55, fill=LIGHT_BLUE, line=BLUE, bold=True)
        add_box(slide, "sequence embedding calibrated for zebrafish", 3.25, 4.6, 2.55, 0.72, fill=WHITE, line=GOLD, size=12, bold=True)
        add_box(slide, "stage-aware or time-conditioned transfer", 6.05, 4.6, 2.55, 0.72, fill=WHITE, line=GREEN, size=12, bold=True)
        add_box(slide, "rollout metrics beyond short displacement", 8.85, 4.6, 2.55, 0.72, fill=WHITE, line=PURPLE, size=12, bold=True)
        add_text(slide, "Main result: modeling more genes is helpful at 10k, especially when UCE state and cross-attention are combined.", 1.05, 6.2, 11.2, 0.45, size=16, bold=True, align=PP_ALIGN.CENTER)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", default="docs/stream_project_slide_deck.pptx")
    args = parser.parse_args()
    create_deck(ROOT / args.output)
    print(f"Wrote {ROOT / args.output}")


if __name__ == "__main__":
    main()
